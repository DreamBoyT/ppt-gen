import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Inches
from io import BytesIO
import re
from langchain_openai import AzureChatOpenAI
from langchain.prompts import PromptTemplate

# Azure OpenAI API details
azure_endpoint = 'https://chat-gpt-a1.openai.azure.com/'
azure_deployment_name = 'DanielChatGPT16k'
azure_api_key = 'c09f91126e51468d88f57cb83a63ee36'
azure_api_version = '2024-05-01-preview'

# Initialize Azure OpenAI LLM
llm = AzureChatOpenAI(
    openai_api_key=azure_api_key,
    api_version=azure_api_version,
    azure_endpoint=azure_endpoint,
    model="gpt-4",
    azure_deployment=azure_deployment_name,
    temperature=0.5
)

def extract_ppt_content(ppt_file):
    prs = Presentation(ppt_file)
    content = []

    for slide_num, slide in enumerate(prs.slides):
        slide_content = {
            "page_number": slide_num + 1,
            "title": slide.shapes.title.text if slide.shapes.title else "",
            "content": [],
            "images": [],
            "tables": [],
            "flow_diagrams": []
        }

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                slide_content["content"].append(shape.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = BytesIO(shape.image.blob)
                slide_content["images"].append(image_stream)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                slide_content["tables"].append(table_data)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Assuming flow diagrams are grouped shapes
                flow_diagram_text = []
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        flow_diagram_text.append(sub_shape.text)
                slide_content["flow_diagrams"].append("\n".join(flow_diagram_text))

        content.append(slide_content)
    return content

def clean_text(text):
    # Remove any control characters or NULL bytes
    return re.sub(r'[\x00-\x1F\x7F]', '', text)

def generate_detailed_explanation(content):
    prompt_template = """
    Given the following slide content, generate a detailed, topic-wise explanation in a point-wise format with topics and subtopics highlighted. Ensure proper spacing, line breaks after every topic and subtopic, indentation, and bolding of key aspects:
    
    Slide Title: {title}
    Slide Content: {content}
    
    Detailed Explanation:
    """
    detailed_content = []
    for slide in content:
        prompt = prompt_template.format(title=slide['title'], content="\n".join(slide['content']))
        response = llm(prompt)
        detailed_content.append(response.content)
    return detailed_content

def create_word_doc(content, detailed_content):
    doc = Document()
    for slide, detailed_text in zip(content, detailed_content):
        title = clean_text(f"Slide {slide['page_number']}: {slide['title']}")
        doc.add_heading(title, level=1)
        doc.add_paragraph(clean_text(detailed_text))
        for paragraph in slide["content"]:
            doc.add_paragraph(clean_text(paragraph))
        for image in slide["images"]:
            try:
                doc.add_picture(image, width=Inches(5.0))
            except Exception as e:
                st.warning(f"Could not add image on slide {slide['page_number']}: {e}")
        for table in slide["tables"]:
            table_doc = doc.add_table(rows=len(table), cols=len(table[0]))
            for i, row in enumerate(table):
                for j, cell in enumerate(row):
                    table_doc.cell(i, j).text = clean_text(cell)
        for flow_diagram in slide["flow_diagrams"]:
            doc.add_paragraph(clean_text(flow_diagram))
    return doc

st.title("PPT to Word Extractor with Detailed Explanations")

uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

if uploaded_file is not None:
    ppt_content = extract_ppt_content(uploaded_file)
    detailed_content = generate_detailed_explanation(ppt_content)
    word_doc = create_word_doc(ppt_content, detailed_content)
    
    buffer = BytesIO()
    word_doc.save(buffer)
    buffer.seek(0)
    
    st.download_button(
        label="Download Word Document",
        data=buffer,
        file_name="extracted_content_with_explanations.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
